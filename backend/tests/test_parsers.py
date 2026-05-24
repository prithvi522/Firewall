import tempfile
from pathlib import Path

import pytest

from app.services.pdf_parser import extract_text_from_file


def test_extract_docx_text_bytes(tmp_path: Path):
    # Create a simple DOCX file using python-docx
    try:
        from docx import Document
    except Exception:
        pytest.skip("python-docx not installed")

    doc = Document()
    doc.add_paragraph("Hello from DOCX test")
    file_path = tmp_path / "test_doc.docx"
    doc.save(str(file_path))

    extracted = extract_text_from_file(file_path, content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document")
    assert "Hello from DOCX test" in extracted


def test_extract_pptx_text_bytes(tmp_path: Path):
    # Create a simple PPTX file using python-pptx
    try:
        from pptx import Presentation
    except Exception:
        pytest.skip("python-pptx not installed")

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    # Use a textbox shape
    from pptx.util import Inches

    left = top = Inches(1)
    width = Inches(6)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Hello from PPTX test"

    file_path = tmp_path / "test_ppt.pptx"
    prs.save(str(file_path))

    extracted = extract_text_from_file(file_path, content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    assert "Hello from PPTX test" in extracted
*** End Patch