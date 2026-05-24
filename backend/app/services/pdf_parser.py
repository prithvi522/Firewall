import re
import subprocess
import tempfile
from pathlib import Path
from shutil import which
from zipfile import BadZipFile, ZipFile
import xml.etree.ElementTree as ET

from PyPDF2 import PdfReader

from app.services.ocr_engine import extract_text_with_ocr

TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".csv",
    ".json",
    ".log",
    ".xml",
    ".yaml",
    ".yml",
    ".toml",
    ".ini",
    ".cfg",
    ".env",
    ".py",
    ".js",
    ".jsx",
    ".ts",
    ".tsx",
    ".dart",
    ".java",
    ".c",
    ".cpp",
    ".cs",
    ".go",
    ".rb",
    ".rs",
    ".php",
    ".sh",
    ".sql",
    ".html",
    ".htm",
    ".css",
}


def extract_text_from_file(file_path: Path, content_type: str | None = None) -> str:
    suffix = file_path.suffix.lower()

    if suffix == ".pdf":
        return extract_text_from_pdf(file_path)

    if suffix == ".docx":
        return extract_text_from_docx(file_path)

    if suffix == ".pptx":
        return extract_text_from_pptx(file_path)

    if suffix in TEXT_EXTENSIONS:
        return read_text_best_effort(file_path)

    if suffix == ".doc":
        return extract_text_from_legacy_office(file_path)

    if suffix == ".ppt":
        return extract_text_from_legacy_office(file_path)

    return read_text_best_effort(file_path)


def read_text_best_effort(file_path: Path) -> str:
    for encoding in ("utf-8", "utf-16", "latin-1"):
        try:
            return file_path.read_text(encoding=encoding, errors="ignore")
        except UnicodeError:
            continue

    return file_path.read_bytes().decode("utf-8", errors="ignore")


def extract_text_from_legacy_office(file_path: Path) -> str:
    converted_text = extract_text_with_libreoffice(file_path)
    if converted_text.strip():
        return converted_text

    return extract_readable_text_from_bytes(file_path.read_bytes())


def extract_text_with_libreoffice(file_path: Path) -> str:
    for executable in ("soffice", "libreoffice"):
        if not which(executable):
            continue

        with tempfile.TemporaryDirectory() as temp_dir:
            output_dir = Path(temp_dir)
            result = subprocess.run(
                [
                    executable,
                    "--headless",
                    "--convert-to",
                    "txt:Text",
                    "--outdir",
                    str(output_dir),
                    str(file_path),
                ],
                capture_output=True,
                text=True,
                timeout=20,
                check=False,
            )

            if result.returncode != 0:
                continue

            converted_file = output_dir / f"{file_path.stem}.txt"
            if converted_file.exists():
                return converted_file.read_text(encoding="utf-8", errors="ignore")

    return ""


def extract_readable_text_from_bytes(content: bytes) -> str:
    printable_chunks = re.findall(rb"[\x20-\x7E]{4,}", content)
    if printable_chunks:
        return "\n".join(chunk.decode("utf-8", errors="ignore") for chunk in printable_chunks)

    return content.decode("utf-8", errors="ignore")


def extract_text_from_pdf(file_path: Path) -> str:
    reader = PdfReader(str(file_path))
    page_text = []

    for page in reader.pages:
        page_text.append(page.extract_text() or "")

    extracted = "\n".join(page_text).strip()
    if extracted:
        return extracted

    # OCR is optional; the OCR service returns an empty string if dependencies are unavailable.
    return extract_text_with_ocr(file_path)


def extract_text_from_docx(file_path: Path) -> str:
    # Prefer python-docx when available for better fidelity; fall back to the zipped XML approach.
    try:
        import docx

        doc = docx.Document(str(file_path))
        paragraphs = [p.text for p in doc.paragraphs if p.text]
        if paragraphs:
            return "\n".join(paragraphs)
    except Exception:
        pass

    try:
        with ZipFile(file_path) as docx:
            xml_content = docx.read("word/document.xml")
    except (BadZipFile, KeyError) as error:
        raise ValueError("Could not read DOCX text from this file.") from error

    root = ET.fromstring(xml_content)
    namespace = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    text_nodes = root.findall(".//w:t", namespace)
    return "\n".join(node.text or "" for node in text_nodes)


def extract_text_from_pptx(file_path: Path) -> str:
    # Prefer python-pptx if installed; otherwise fall back to XML-in-zip extraction.
    try:
        from pptx import Presentation

        prs = Presentation(str(file_path))
        slides = []
        for slide in prs.slides:
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        texts.append(text)
            if texts:
                slides.append("\n".join(texts))
        if slides:
            return "\n\n".join(slides)
    except Exception:
        pass

    try:
        with ZipFile(file_path) as pptx:
            slide_names = sorted(
                name for name in pptx.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")
            )

            if not slide_names:
                raise KeyError("No slides found in this file.")

            namespace = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
            slide_texts: list[str] = []

            for slide_name in slide_names:
                root = ET.fromstring(pptx.read(slide_name))
                text_nodes = root.findall(".//a:t", namespace)
                slide_text = "\n".join(node.text or "" for node in text_nodes).strip()
                if slide_text:
                    slide_texts.append(slide_text)

            return "\n\n".join(slide_texts)
    except (BadZipFile, KeyError, ET.ParseError) as error:
        raise ValueError("Could not read PPTX text from this file.") from error
